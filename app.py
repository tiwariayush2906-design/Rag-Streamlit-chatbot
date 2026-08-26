import streamlit as st
from datetime import datetime
from zoneinfo import ZoneInfo
import os
import uuid
import tempfile
import random
import io

from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from dotenv import load_dotenv
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from openpyxl import Workbook

load_dotenv()

# Streamlit secrets ya environment variable se API key fetch karna
groq_api_key = st.secrets.get("GROQ_API_KEY", os.getenv("GROQ_API_KEY"))

st.set_page_config(page_title="RAG Chatbot", page_icon="🤖", layout="centered", initial_sidebar_state="expanded")

st.markdown("""
<style>
#MainMenu, footer {visibility: hidden;}

header {
    visibility: visible !important;
    background: transparent !important;
}
header [data-testid="stToolbar"] {
    visibility: hidden;
}

.block-container {
    padding-top: 3rem;
    max-width: 750px;
}

h1 {
    font-family: Georgia, 'Times New Roman', serif;
    font-weight: 400;
    font-size: 2.6rem;
    text-align: center;
    margin-bottom: 0.2rem;
}

.subtitle {
    text-align: center;
    color: #888;
    font-size: 1rem;
    margin-bottom: 2rem;
}

[data-testid="stChatInput"] {
    border-radius: 25px;
}

.loading-dots {
    display: flex;
    align-items: center;
    gap: 6px;
    padding: 10px 0;
}

.loading-dots span {
    width: 8px;
    height: 8px;
    border-radius: 50%;
    background-color: #d97757;
    animation: bounce 1.4s infinite ease-in-out both;
}

.loading-dots span:nth-child(1) { animation-delay: -0.32s; }
.loading-dots span:nth-child(2) { animation-delay: -0.16s; }

@keyframes bounce {
    0%, 80%, 100% { transform: scale(0.6); opacity: 1; }
    40% { transform: scale(1); opacity: 1; }
}

[data-testid="stSidebar"] button {
    text-align: left !important;
    justify-content: flex-start !important;
}

.sidebar-empty {
    color: #999;
    font-size: 0.85rem;
    padding: 10px 4px;
}
</style>
""", unsafe_allow_html=True)

# =========================================================
# Multi-chat session state
# =========================================================
if "chat_sessions" not in st.session_state:
    st.session_state.chat_sessions = {}
if "current_chat_id" not in st.session_state:
    st.session_state.current_chat_id = None
if "response_language" not in st.session_state:
    st.session_state.response_language = "english"


def new_chat_dict():
    return {
        "title": "New Chat",
        "messages": [],
        "attached_files": [],
        "vectorstore": None,
        "retriever": None,
    }


def ensure_active_chat():
    if st.session_state.current_chat_id is None or \
       st.session_state.current_chat_id not in st.session_state.chat_sessions:
        chat_id = str(uuid.uuid4())
        st.session_state.chat_sessions[chat_id] = new_chat_dict()
        st.session_state.current_chat_id = chat_id


def current_chat():
    ensure_active_chat()
    return st.session_state.chat_sessions[st.session_state.current_chat_id]


def start_new_chat():
    chat_id = str(uuid.uuid4())
    st.session_state.chat_sessions[chat_id] = new_chat_dict()
    st.session_state.current_chat_id = chat_id


def switch_chat(chat_id):
    st.session_state.current_chat_id = chat_id


ensure_active_chat()

# =========================================================
# Sidebar — New Chat + Search History
# =========================================================
with st.sidebar:
    st.markdown("### 🤖 RAG Chatbot")

    if st.button("➕ New Chat", use_container_width=True, type="primary"):
        start_new_chat()
        st.rerun()

    st.markdown("---")
    st.markdown("#### 🕘 History")

    search_query = st.text_input(
        "Search history",
        label_visibility="collapsed",
        placeholder="🔍 Search chats..."
    )

    all_chat_ids = list(st.session_state.chat_sessions.keys())[::-1]

    visible_chats = [
        cid for cid in all_chat_ids
        if st.session_state.chat_sessions[cid]["messages"]
    ]

    if search_query:
        visible_chats = [
            cid for cid in visible_chats
            if search_query.lower() in st.session_state.chat_sessions[cid]["title"].lower()
        ]

    if not visible_chats:
        st.markdown('<p class="sidebar-empty">Abhi koi purani chat nahi hai.</p>', unsafe_allow_html=True)
    else:
        for cid in visible_chats:
            chat_item = st.session_state.chat_sessions[cid]
            is_active = cid == st.session_state.current_chat_id
            label = f"{'💬 ' if not is_active else '▶️ '}{chat_item['title']}"
            if st.button(label, key=f"hist_{cid}", use_container_width=True):
                switch_chat(cid)
                st.rerun()

# =========================================================
# Dynamic Timezone Based Greeting (Asia/Kolkata IST)
# =========================================================
current_hour = datetime.now(ZoneInfo("Asia/Kolkata")).hour

if 5 <= current_hour < 12:
    greeting = "Good morning"
elif 12 <= current_hour < 17:
    greeting = "Good afternoon"
elif 17 <= current_hour < 21:
    greeting = "Good evening"
else:
    greeting = "Good night"

st.markdown(f"<h1>🌟 {greeting}, dost!</h1>", unsafe_allow_html=True)
st.markdown('<p class="subtitle">How can I help you today?</p>', unsafe_allow_html=True)

LANGUAGE_MAP = {
    "english": "Poora answer sirf English me do.",
    "hindi": "Poora answer sirf Hindi (Devanagari script) me do.",
    "hinglish": "Poora answer Hinglish (Hindi-English mix, Roman script) me do."
}

RAG_PROMPT_TEMPLATE = """Tum ek intelligent, highly accurate assistant ho jo attached document ke content ke saath kaam karte ho.

Neeche diya gaya context tumhari attached file(s) ka content hai.

Recent Chat History:
{chat_history}

RESPONSE LENGTH & DETAIL RULES:
1. Normal Question Par: Exact, compact aur direct answer do.
2. Explanation Command Par: Jab user keh-e "detail me samjhao", "explain karo", ya "vistaar se batao", tabhi full detailed step-by-step response do.

ACCURACY RULES:
- Point-to-point exact facts document context se do.
- Agar context me information na mile (jaise Competitive exams, Government services, Indian laws, graduation subjects, Indian leaders, politics, celebrity bios, formulas, English grammar, coding/SQL, human biology/psychology, defense, ya medical), toh apni internal knowledge base se accurate answer do.

LANGUAGE RULES:
- {language_instruction}
- Answer ki shuruaat hamesha ek English word/phrase se karo, phir baaki jawab us language me.

Context: {context}

User ka message: {question}

Response:"""

GENERAL_SYSTEM_PROMPT = """Tum ek universal, highly knowledgeable, aur ultra-accurate AI assistant ho.

EXPERT KNOWLEDGE AREAS:
1. COMPETITIVE EXAMS & GOVERNMENT SERVICES: Complete syllabi, exam patterns, eligibility, strategy, previous year trends, officer ranks, pay structures, and job profiles for UPSC (CSE, CAPF, CDS, NDA, ESE), State PSCs, SSC (CGL, CHSL, CPO, JE), Banking & Finance (IBPS, SBI PO/Clerk, RBI Grade B, SEBI, NABARD), Railways (RRB NTPC, JE, ALP), Defence Services, GATE, Judicial Services (PCS-J), and NTA NET/JRF.
2. INDIAN LAW, RULES & CONSTITUTION: Complete Indian legal system, Constitution of India (Articles, Fundamental Rights, Writs, Amendments, Landmark Judgments), Bharatiya Nyaya Sanhita (BNS), Bharatiya Nagarik Suraksha Sanhita (BNSS), Bharatiya Sakshya Adhiniyam (BSA), Contract Act, Company Law, Cyber Laws, Tax Laws (GST & Income Tax), Family Laws, and Consumer Protection.
3. GRADUATION & HIGHER EDUCATION SUBJECTS: Deep academic knowledge across B.Tech/BCA/MCA (Computer Science, Data Structures, Engineering Math), B.Sc (Physics, Chemistry, Advanced Mathematics, Bio-tech), B.Com/BBA/MBA (Accounting, Economics, Business Law, Finance), and B.A. (Political Science, Sociology, History, Public Administration, Literature).
4. LEADERS, HISTORY & POLITICS: Current & historical leaders of India & world, biographies, political parties (ideologies, history, structures), freedom movement, Indian dynasties, and geopolitical developments.
5. CELEBRITIES & BIOGRAPHIES: Detailed biographies, career timelines, achievements of famous personalities across movies, sports, arts, science, and business.
6. FORMULAS & SCIENTIFIC LAWS: Comprehensive mathematical formulas, Physics principles/equations, Chemical reactions, and engineering calculations.
7. ENGLISH GRAMMAR & LANGUAGE: Syntax, parts of speech, tense usage, voice, narration, clause structure, idioms, advanced vocabulary, and writing techniques.
8. CODING & SQL: Python, JavaScript, C++, Java, Web Development, Database Management, SQL queries, Data Structures, Algorithms, and debugging.
9. HUMAN ANATOMY & PSYCHOLOGY: Human biological systems, organ functions, medical physiology, cognitive psychology, emotional intelligence, and human behavior.
10. MILITARY, DEFENSE & MEDICINE: Armed forces, strategic weapons, medicines, formulations, ointments/tubes, and emergency first-aid.

Aaj ki date aur time: {current_datetime}
Agar user "aaj", "aaj ka din", ya current date se related poochta hai, tabhi upar di gayi date use karo.

Recent Chat History:
{chat_history}

RESPONSE STYLE RULES:
- Fast & Direct: Normal query par fast, point-to-point accurate answer do.
- Deep Explanation On Request: "Explain karo", "detail me samjhao", "vistaar se batao" kahne par hi thorough detailed explanation do.

LANGUAGE RULES:
- {language_instruction}
- Answer ki shuruaat hamesha ek English word/phrase se karo, phir baaki jawab us language me."""

VERIFY_PROMPT = """Tum ek strict fact-checker aur editor ho. Neeche ek AI dwara diya gaya DRAFT ANSWER hai jo verify karna hai.

Original question: {question}
{context_block}
Draft Answer:
{draft_answer}

Tumhara kaam:
1. Check karo ki koi Competitive Exam rule, Government Service details, Legal fault (Indian Laws/Articles), Academic concept error, Factual inaccuracy, Political mistake, Formula flaw, ya Code bug na ho.
2. Normal queries ke liye compact short response aur explicitly detail mange jane par thorough detailed format maintain rakho.
3. Silent fix karke clean, highly accurate final output do.

Sirf clean final answer do - koi meta-commentary mat likho."""


@st.cache_resource
def get_embeddings():
    return HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")


@st.cache_resource
def get_llm():
    return ChatGroq(
        model="llama-3.3-70b-versatile",
        temperature=0,
        groq_api_key=groq_api_key
    )


def verify_and_correct(question, draft_answer, source_docs=None):
    llm = get_llm()

    if source_docs:
        context_text = "\n\n".join(doc.page_content for doc in source_docs)
        context_block = f"\nReference context (document se):\n{context_text}\n"
    else:
        context_block = ""

    verify_prompt = VERIFY_PROMPT.format(
        question=question,
        context_block=context_block,
        draft_answer=draft_answer
    )

    try:
        response = llm.invoke([("human", verify_prompt)])
        corrected = response.content.strip()
        return corrected if corrected else draft_answer
    except Exception:
        return draft_answer


def load_file(file_path, file_extension):
    if file_extension == "pdf":
        return PyPDFLoader(file_path).load()
    elif file_extension == "txt":
        return TextLoader(file_path, encoding="utf-8").load()
    elif file_extension == "docx":
        return Docx2txtLoader(file_path).load()
    return []


def process_file(uploaded_file):
    chat = current_chat()

    file_extension = uploaded_file.name.split(".")[-1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_extension}") as tmp_file:
        tmp_file.write(uploaded_file.read())
        tmp_path = tmp_file.name

    documents = load_file(tmp_path, file_extension)
    splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=200)
    chunks = splitter.split_documents(documents)

    embeddings = get_embeddings()
    collection_name = f"chat_{st.session_state.current_chat_id}".replace("-", "_")
    
    if chat["vectorstore"] is None:
        chat["vectorstore"] = Chroma.from_documents(
            documents=chunks, embedding=embeddings, collection_name=collection_name
        )
    else:
        chat["vectorstore"].add_documents(chunks)

    os.unlink(tmp_path)
    chat["retriever"] = chat["vectorstore"].as_retriever(
        search_type="mmr",
        search_kwargs={"k": 5, "fetch_k": 10, "lambda_mult": 0.7}
    )


def detect_language_command(text):
    t = text.lower()
    if any(x in t for x in ["hindi me do", "hindi mein do", "ab hindi", "hindi me bolo"]):
        return "hindi"
    if any(x in t for x in ["hinglish me do", "hinglish mein do", "ab hinglish"]):
        return "hinglish"
    if any(x in t for x in ["english me do", "ab english", "english me bolo"]):
        return "english"
    return None


def detect_file_format(text):
    t = text.lower()
    if "pdf" in t:
        return "pdf"
    if any(x in t for x in ["word", "docx", "doc file"]):
        return "docx"
    if any(x in t for x in ["excel", "xlsx", "spreadsheet"]):
        return "xlsx"
    if any(x in t for x in ["ppt", "pptx", "powerpoint", "presentation", "slides"]):
        return "pptx"
    return None


ANSWER_COLORS = [
    "#1e1e1e",
    "#1b3a4b",
    "#2d1b4e",
    "#1b4d3e",
    "#4d1b1b",
    "#4d3a1b",
    "#1b4d4d",
    "#3a1b4d",
]


def render_answer(text):
    color = random.choice(ANSWER_COLORS)
    st.markdown(
        f"""<div style="background-color:{color};color:#f5f5f5;padding:16px;
        border-radius:10px;line-height:1.6;font-size:16px;white-space:pre-wrap;">{text}</div>""",
        unsafe_allow_html=True
    )


def generate_pdf(text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    for line in text.split("\n"):
        safe_line = line.encode("latin-1", errors="replace").decode("latin-1")
        pdf.multi_cell(0, 8, safe_line)
    return bytes(pdf.output())


def generate_docx(text):
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def generate_pptx(text):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Generated Content"
    body = slide.placeholders[1].text_frame
    lines = [l for l in text.split("\n") if l.strip()]
    if lines:
        body.text = lines[0]
        for line in lines[1:]:
            p = body.add_paragraph()
            p.text = line
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_xlsx(text):
    wb = Workbook()
    ws = wb.active
    for i, line in enumerate(text.split("\n"), start=1):
        ws.cell(row=i, column=1, value=line)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


FILE_GENERATORS = {
    "pdf": (generate_pdf, "application/pdf", "output.pdf"),
    "docx": (generate_docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "output.docx"),
    "pptx": (generate_pptx, "application/vnd.openxmlformats-officedocument.presentationml.presentation", "output.pptx"),
    "xlsx": (generate_xlsx, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "output.xlsx"),
}

chat = current_chat()

if chat["attached_files"]:
    st.caption("📎 Attached: " + ", ".join(chat["attached_files"]))

for msg in chat["messages"]:
    with st.chat_message(msg["role"]):
        if msg["role"] == "assistant":
            render_answer(msg["content"])
        else:
            st.write(msg["content"])

chat_input_data = st.chat_input(
    "Poochiye kuch bhi, ya seedha file attach kar dijiye... ✨",
    accept_file=True,
    file_type=["pdf", "txt", "docx"]
)

user_query = None
newly_uploaded_file = None

if chat_input_data:
    user_query = chat_input_data.text
    if chat_input_data.files:
        newly_uploaded_file = chat_input_data.files[0]

if newly_uploaded_file is not None:
    with st.spinner(f"'{newly_uploaded_file.name}' process ho raha hai..."):
        process_file(newly_uploaded_file)
        chat["attached_files"].append(newly_uploaded_file.name)
    st.success(f"✅ '{newly_uploaded_file.name}' ready hai!")

if user_query:
    if chat["title"] == "New Chat":
        chat["title"] = (user_query[:30] + "…") if len(user_query) > 30 else user_query

    lang_cmd = detect_language_command(user_query)

    if lang_cmd:
        st.session_state.response_language = lang_cmd
        chat["messages"].append({"role": "user", "content": user_query})
        with st.chat_message("user"):
            st.write(user_query)
        confirm_msg = {
            "hindi": "ठीक है, अब मैं हिंदी में जवाब दूंगा।",
            "hinglish": "Theek hai, ab main Hinglish me jawab dunga.",
            "english": "Sure, I'll answer in English from now on."
        }[lang_cmd]
        with st.chat_message("assistant"):
            render_answer(confirm_msg)
        chat["messages"].append({"role": "assistant", "content": confirm_msg})

    else:
        chat["messages"].append({"role": "user", "content": user_query})
        with st.chat_message("user"):
            st.write(user_query)

        with st.chat_message("assistant"):
            loading_placeholder = st.empty()
            loading_placeholder.markdown(
                '<div class="loading-dots"><span></span><span></span><span></span></div>',
                unsafe_allow_html=True
            )

            language_instruction = LANGUAGE_MAP[st.session_state.response_language]
            llm = get_llm()

            # Format recent chat history (last 6 messages)
            recent_messages = chat["messages"][-6:]
            history_text = ""
            for m in recent_messages:
                role_label = "User" if m["role"] == "user" else "Assistant"
                history_text += f"{role_label}: {m['content']}\n"

            if chat["retriever"] is not None:
                source_docs = chat["retriever"].invoke(user_query)
                context_text = "\n\n".join(doc.page_content for doc in source_docs)

                rag_prompt = ChatPromptTemplate.from_template(RAG_PROMPT_TEMPLATE)
                chain = rag_prompt | llm | StrOutputParser()

                draft_answer = chain.invoke({
                    "context": context_text,
                    "question": user_query,
                    "language_instruction": language_instruction,
                    "chat_history": history_text
                })
            else:
                current_datetime = datetime.now(ZoneInfo("Asia/Kolkata")).strftime("%A, %d %B %Y, %I:%M %p")
                
                system_msg = GENERAL_SYSTEM_PROMPT.format(
                    language_instruction=language_instruction,
                    current_datetime=current_datetime,
                    chat_history=history_text
                )
                response = llm.invoke([
                    ("system", system_msg),
                    ("human", user_query)
                ])
                draft_answer = response.content
                source_docs = None

            answer = verify_and_correct(user_query, draft_answer, source_docs)

            loading_placeholder.empty()
            render_answer(answer)

            file_format = detect_file_format(user_query)
            if file_format:
                gen_func, mime, filename = FILE_GENERATORS[file_format]
                file_bytes = gen_func(answer)
                st.download_button(
                    label=f"📥 Download as {file_format.upper()}",
                    data=file_bytes,
                    file_name=filename,
                    mime=mime
                )

        chat["messages"].append({"role": "assistant", "content": answer})